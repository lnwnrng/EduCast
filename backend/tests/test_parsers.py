"""文档解析器单元测试。

覆盖:
  - PPTX 解析（python-pptx 生成的临时文件）
  - PDF 解析（pdfplumber）
  - Markdown 解析
  - 纯文本解析
  - 不支持的文件类型
  - 文件不存在
  - 空文件
  - IR 构建: slide → scene 映射
"""

import os
import tempfile
from pathlib import Path

import pytest
import pytest_asyncio

from app.exceptions import ParseException
from app.ir.schema import CourseIR, SceneType
from app.pipeline.parser import DocumentParser, ParsedSlide


# ── 测试 fixtures ────────────────────────────────────────


@pytest.fixture
def parser(tmp_path: Path) -> DocumentParser:
    """创建解析器实例（使用临时存储目录）。"""
    return DocumentParser(storage_root=str(tmp_path / "storage"))


@pytest.fixture
def sample_pptx(tmp_path: Path) -> str:
    """创建一个包含 3 页的测试 PPTX 文件。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # 第 1 页: 标题页
    title_layout = prs.slide_layouts[0]  # Title Slide
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "高等数学导论"
    slide1.placeholders[1].text = "第一学期课程"

    # 第 2 页: 内容页
    content_layout = prs.slide_layouts[1]  # Title and Content
    slide2 = prs.slides.add_slide(content_layout)
    slide2.shapes.title.text = "导数的定义"
    slide2.placeholders[1].text = "导数是函数的变化率\n表示瞬时速度"
    # 添加备注
    notes_slide = slide2.notes_slide
    notes_slide.notes_text_frame.text = "讲解导数的基本概念，注意强调极限定义。"

    # 第 3 页: 内容页
    slide3 = prs.slides.add_slide(content_layout)
    slide3.shapes.title.text = "导数的计算"
    slide3.placeholders[1].text = "基本公式：(x^n)' = nx^{n-1}"

    file_path = str(tmp_path / "test_math.pptx")
    prs.save(file_path)
    return file_path


@pytest.fixture
def sample_pptx_single_slide(tmp_path: Path) -> str:
    """创建一个只有 1 页的 PPTX 文件。"""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "单页测试"
    slide.placeholders[1].text = "这是唯一的一页内容"

    file_path = str(tmp_path / "single.pptx")
    prs.save(file_path)
    return file_path


@pytest.fixture
def sample_markdown(tmp_path: Path) -> str:
    """创建一个测试 Markdown 文件。"""
    content = """# 数据结构基础

## 数组

数组是最基本的数据结构，支持随机访问。

时间复杂度：
- 访问: O(1)
- 搜索: O(n)

## 链表

链表是一种线性数据结构。

每个节点包含数据和指针。

## 栈与队列

### 栈

后进先出 (LIFO) 数据结构。

### 队列

先进先出 (FIFO) 数据结构。
"""
    file_path = str(tmp_path / "data_structures.md")
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(content)
    return file_path


@pytest.fixture
def sample_markdown_no_headings(tmp_path: Path) -> str:
    """创建一个没有标题的 Markdown 文件。"""
    content = """这是一段纯文本内容。

没有任何标题标记。

只有段落文本。
"""
    file_path = str(tmp_path / "no_headings.md")
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(content)
    return file_path


@pytest.fixture
def sample_text(tmp_path: Path) -> str:
    """创建一个测试纯文本文件。"""
    content = """Python 编程基础

Python 是一种高级编程语言，以简洁易读著称。

变量和数据类型是编程的基础。Python 支持整数、浮点数、字符串等基本数据类型。

函数是组织代码的基本单元。使用 def 关键字定义函数。
"""
    file_path = str(tmp_path / "python_basics.txt")
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(content)
    return file_path


@pytest.fixture
def sample_empty_text(tmp_path: Path) -> str:
    """创建一个空文本文件。"""
    file_path = str(tmp_path / "empty.txt")
    with open(file_path, "w", encoding="utf-8") as f:
        f.write("")
    return file_path


@pytest.fixture
def sample_pdf(tmp_path: Path) -> str:
    """创建一个简单的测试 PDF 文件。

    使用 fpdf2 或 reportlab 生成。
    如果不可用，创建一个最小的 PDF。
    """
    file_path = str(tmp_path / "test_doc.pdf")

    # 使用最基本的 PDF 格式手动创建
    # 这是一个包含单页文本的最小 PDF
    pdf_content = b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj
4 0 obj
<< /Length 44 >>
stream
BT /F1 12 Tf 100 700 Td (Test PDF Page) Tj ET
endstream
endobj
5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj
xref
0 6
0000000000 65535 f 
0000000009 00000 n 
0000000058 00000 n 
0000000115 00000 n 
0000000266 00000 n 
0000000360 00000 n 
trailer
<< /Size 6 /Root 1 0 R >>
startxref
441
%%EOF"""

    with open(file_path, "wb") as f:
        f.write(pdf_content)
    return file_path


# ── PPTX 解析测试 ────────────────────────────────────────


class TestPPTXParser:
    """PPTX 文件解析测试。"""

    @pytest.mark.asyncio
    async def test_parse_pptx_basic(
        self, parser: DocumentParser, sample_pptx: str
    ) -> None:
        """解析包含 3 页的 PPTX 文件。"""
        ir = await parser.parse(sample_pptx, ".pptx")

        assert isinstance(ir, CourseIR)
        assert ir.title == "test_math"
        assert len(ir.chapters) >= 1

        # 至少应该有场景
        total_scenes = sum(
            len(kp.scenes)
            for ch in ir.chapters
            for kp in ch.knowledge_points
        )
        assert total_scenes >= 1

    @pytest.mark.asyncio
    async def test_parse_pptx_extracts_notes(
        self, parser: DocumentParser, sample_pptx: str
    ) -> None:
        """PPTX 解析应提取备注作为旁白。"""
        ir = await parser.parse(sample_pptx, ".pptx")

        # 遍历所有场景，检查是否有包含备注的旁白
        all_narrations = [
            scene.narration_text
            for ch in ir.chapters
            for kp in ch.knowledge_points
            for scene in kp.scenes
        ]
        # 第 2 页有备注
        has_notes = any("导数" in n or "极限" in n for n in all_narrations)
        assert has_notes, f"备注未被提取到旁白中: {all_narrations}"

    @pytest.mark.asyncio
    async def test_parse_pptx_scene_type_is_slide(
        self, parser: DocumentParser, sample_pptx: str
    ) -> None:
        """PPTX 解析的所有分镜应为 SLIDE 类型。"""
        ir = await parser.parse(sample_pptx, ".pptx")

        for ch in ir.chapters:
            for kp in ch.knowledge_points:
                for scene in kp.scenes:
                    assert scene.scene_type == SceneType.SLIDE

    @pytest.mark.asyncio
    async def test_parse_pptx_single_slide(
        self, parser: DocumentParser, sample_pptx_single_slide: str
    ) -> None:
        """解析单页 PPTX。"""
        ir = await parser.parse(sample_pptx_single_slide, ".pptx")

        assert isinstance(ir, CourseIR)
        assert len(ir.chapters) >= 1
        total_scenes = sum(
            len(kp.scenes)
            for ch in ir.chapters
            for kp in ch.knowledge_points
        )
        assert total_scenes >= 1

    @pytest.mark.asyncio
    async def test_parse_pptx_scenes_have_source_page(
        self, parser: DocumentParser, sample_pptx: str
    ) -> None:
        """分镜应记录来源页码。"""
        ir = await parser.parse(sample_pptx, ".pptx")

        for ch in ir.chapters:
            for kp in ch.knowledge_points:
                for scene in kp.scenes:
                    assert scene.source_page is not None
                    assert scene.source_page >= 1

    @pytest.mark.asyncio
    async def test_parse_pptx_scenes_have_slide_ref(
        self, parser: DocumentParser, sample_pptx: str
    ) -> None:
        """SLIDE 类型分镜应有 slide_ref。"""
        ir = await parser.parse(sample_pptx, ".pptx")

        for ch in ir.chapters:
            for kp in ch.knowledge_points:
                for scene in kp.scenes:
                    assert scene.visual_spec.slide_ref


# ── Markdown 解析测试 ────────────────────────────────────


class TestMarkdownParser:
    """Markdown 文件解析测试。"""

    @pytest.mark.asyncio
    async def test_parse_markdown_basic(
        self, parser: DocumentParser, sample_markdown: str
    ) -> None:
        """解析 Markdown 文件。"""
        ir = await parser.parse(sample_markdown, ".md")

        assert isinstance(ir, CourseIR)
        assert ir.title == "data_structures"
        assert len(ir.chapters) >= 1

    @pytest.mark.asyncio
    async def test_parse_markdown_chapter_detection(
        self, parser: DocumentParser, sample_markdown: str
    ) -> None:
        """Markdown # 标题应映射为章节。"""
        ir = await parser.parse(sample_markdown, ".md")

        # 应有 1 个 # 标题 → 1 个章节
        assert len(ir.chapters) >= 1
        assert "数据结构" in ir.chapters[0].title

    @pytest.mark.asyncio
    async def test_parse_markdown_knowledge_points(
        self, parser: DocumentParser, sample_markdown: str
    ) -> None:
        """Markdown ## 标题应映射为知识点。"""
        ir = await parser.parse(sample_markdown, ".md")

        kp_titles = [
            kp.title
            for ch in ir.chapters
            for kp in ch.knowledge_points
        ]
        assert any("数组" in t for t in kp_titles)
        assert any("链表" in t for t in kp_titles)

    @pytest.mark.asyncio
    async def test_parse_markdown_body_to_narration(
        self, parser: DocumentParser, sample_markdown: str
    ) -> None:
        """Markdown 正文应映射为分镜旁白。"""
        ir = await parser.parse(sample_markdown, ".md")

        all_narrations = [
            scene.narration_text
            for ch in ir.chapters
            for kp in ch.knowledge_points
            for scene in kp.scenes
        ]
        has_content = any("数据结构" in n or "数组" in n for n in all_narrations)
        assert has_content

    @pytest.mark.asyncio
    async def test_parse_markdown_no_headings(
        self, parser: DocumentParser, sample_markdown_no_headings: str
    ) -> None:
        """无标题 Markdown 应作为纯文本处理。"""
        ir = await parser.parse(sample_markdown_no_headings, ".md")

        assert isinstance(ir, CourseIR)
        assert len(ir.chapters) >= 1


# ── 纯文本解析测试 ───────────────────────────────────────


class TestTextParser:
    """纯文本解析测试。"""

    @pytest.mark.asyncio
    async def test_parse_text_basic(
        self, parser: DocumentParser, sample_text: str
    ) -> None:
        """解析纯文本文件。"""
        ir = await parser.parse(sample_text, ".txt")

        assert isinstance(ir, CourseIR)
        assert ir.title == "python_basics"
        assert len(ir.chapters) == 1

    @pytest.mark.asyncio
    async def test_parse_text_paragraphs_to_scenes(
        self, parser: DocumentParser, sample_text: str
    ) -> None:
        """纯文本段落应映射为分镜。"""
        ir = await parser.parse(sample_text, ".txt")

        total_scenes = sum(
            len(kp.scenes)
            for ch in ir.chapters
            for kp in ch.knowledge_points
        )
        # 多个段落 → 多个场景
        assert total_scenes >= 2

    @pytest.mark.asyncio
    async def test_parse_text_scene_order(
        self, parser: DocumentParser, sample_text: str
    ) -> None:
        """分镜序号应连续。"""
        ir = await parser.parse(sample_text, ".txt")

        for ch in ir.chapters:
            for kp in ch.knowledge_points:
                for idx, scene in enumerate(kp.scenes):
                    assert scene.order == idx + 1


# ── PDF 解析测试 ─────────────────────────────────────────


class TestPDFParser:
    """PDF 文件解析测试。"""

    @pytest.mark.asyncio
    async def test_parse_pdf_basic(
        self, parser: DocumentParser, sample_pdf: str
    ) -> None:
        """解析 PDF 文件。"""
        ir = await parser.parse(sample_pdf, ".pdf")

        assert isinstance(ir, CourseIR)
        assert len(ir.chapters) >= 1


# ── 错误处理测试 ─────────────────────────────────────────


class TestErrorHandling:
    """解析器错误处理测试。"""

    @pytest.mark.asyncio
    async def test_unsupported_extension(
        self, parser: DocumentParser
    ) -> None:
        """不支持的文件类型应抛出 ParseException。"""
        with pytest.raises(ParseException, match="不支持的文件类型"):
            await parser.parse("test.xyz", ".xyz")

    @pytest.mark.asyncio
    async def test_file_not_found(
        self, parser: DocumentParser
    ) -> None:
        """文件不存在应抛出 ParseException。"""
        with pytest.raises(ParseException, match="文件不存在"):
            await parser.parse("/nonexistent/file.pptx", ".pptx")

    @pytest.mark.asyncio
    async def test_empty_text_file(
        self, parser: DocumentParser, sample_empty_text: str
    ) -> None:
        """空文本文件应抛出 ParseException。"""
        with pytest.raises(ParseException):
            await parser.parse(sample_empty_text, ".txt")


# ── ParsedSlide 测试 ────────────────────────────────────


class TestParsedSlide:
    """解析中间结构测试。"""

    def test_has_content_with_title(self) -> None:
        """有标题即算有内容。"""
        slide = ParsedSlide(page_number=1, title="标题")
        assert slide.has_content is True

    def test_has_content_with_body(self) -> None:
        """有正文即算有内容。"""
        slide = ParsedSlide(page_number=1, body_text="正文")
        assert slide.has_content is True

    def test_has_content_with_notes(self) -> None:
        """有备注即算有内容。"""
        slide = ParsedSlide(page_number=1, notes="备注")
        assert slide.has_content is True

    def test_no_content(self) -> None:
        """空 slide 无内容。"""
        slide = ParsedSlide(page_number=1)
        assert slide.has_content is False


# ── IR 构建逻辑测试 ─────────────────────────────────────


class TestIRBuildingLogic:
    """IR 构建逻辑测试。"""

    @pytest.mark.asyncio
    async def test_ir_has_four_layers(
        self, parser: DocumentParser, sample_pptx: str
    ) -> None:
        """IR 应包含四层结构: Course → Chapter → KP → Scene。"""
        ir = await parser.parse(sample_pptx, ".pptx")

        assert ir.title  # Course
        assert ir.chapters  # Chapters
        for ch in ir.chapters:
            assert ch.title  # Chapter title
            assert ch.knowledge_points  # KPs
            for kp in ch.knowledge_points:
                assert kp.title  # KP title
                assert kp.scenes  # Scenes

    @pytest.mark.asyncio
    async def test_ir_scenes_are_slide_type(
        self, parser: DocumentParser, sample_pptx: str
    ) -> None:
        """解析产出的分镜应全部为 SLIDE 类型。"""
        ir = await parser.parse(sample_pptx, ".pptx")

        for ch in ir.chapters:
            for kp in ch.knowledge_points:
                for scene in kp.scenes:
                    assert scene.scene_type == SceneType.SLIDE

    @pytest.mark.asyncio
    async def test_ir_json_serializable(
        self, parser: DocumentParser, sample_pptx: str
    ) -> None:
        """IR 应可序列化为 JSON。"""
        ir = await parser.parse(sample_pptx, ".pptx")

        json_str = ir.model_dump_json()
        assert json_str

        # 反序列化验证
        ir2 = CourseIR.model_validate_json(json_str)
        assert ir2.title == ir.title
        assert len(ir2.chapters) == len(ir.chapters)
