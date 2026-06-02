"""上传 API 集成测试。

覆盖:
  - 文件上传成功 → 创建 Project + Task
  - 文件类型校验
  - 空文件校验
  - 上传后 IR 生成
  - 脚本 API (get_script, update_script)
"""

import io

import pytest
from httpx import AsyncClient
from sqlalchemy.ext.asyncio import AsyncSession

PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument." "presentationml.presentation"
)


# ── 辅助函数 ─────────────────────────────────────────────


def _make_pptx_bytes() -> bytes:
    """生成一个最小的 PPTX 文件字节。"""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "测试上传"
    slide.placeholders[1].text = "这是测试内容"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_md_bytes() -> bytes:
    """生成一个 Markdown 文件字节。"""
    content = """# 测试课程

## 第一节

这是第一节的内容。

## 第二节

这是第二节的内容。
"""
    return content.encode("utf-8")


def _make_txt_bytes() -> bytes:
    """生成一个纯文本文件字节。"""
    return "纯文本测试内容\n\n这是第二段。".encode()


# ── 上传 API 测试 ────────────────────────────────────────


class TestUploadAPI:
    """文件上传 API 测试。"""

    @pytest.mark.asyncio
    async def test_upload_pptx_success(self, client: AsyncClient) -> None:
        """上传 PPTX 文件成功。"""
        pptx_bytes = _make_pptx_bytes()

        resp = await client.post(
            "/api/v1/upload/document",
            files={
                "file": ("test.pptx", pptx_bytes, PPTX_MIME),
            },
        )

        assert resp.status_code == 200
        data = resp.json()
        assert data["message"] == "文件上传成功，解析任务已创建"
        assert "project_id" in data["data"]
        assert "task_id" in data["data"]
        assert data["data"]["file_type"] == ".pptx"

    @pytest.mark.asyncio
    async def test_upload_markdown_success(self, client: AsyncClient) -> None:
        """上传 Markdown 文件成功。"""
        md_bytes = _make_md_bytes()

        resp = await client.post(
            "/api/v1/upload/document",
            files={
                "file": ("course.md", md_bytes, "text/markdown"),
            },
        )

        assert resp.status_code == 200
        data = resp.json()
        assert data["data"]["file_type"] == ".md"
        assert data["data"]["project_id"]

    @pytest.mark.asyncio
    async def test_upload_txt_success(self, client: AsyncClient) -> None:
        """上传纯文本文件成功。"""
        txt_bytes = _make_txt_bytes()

        resp = await client.post(
            "/api/v1/upload/document",
            files={
                "file": ("notes.txt", txt_bytes, "text/plain"),
            },
        )

        assert resp.status_code == 200
        data = resp.json()
        assert data["data"]["file_type"] == ".txt"

    @pytest.mark.asyncio
    async def test_upload_unsupported_type(self, client: AsyncClient) -> None:
        """上传不支持的文件类型应返回 422 错误。"""
        resp = await client.post(
            "/api/v1/upload/document",
            files={
                "file": ("test.xyz", b"dummy content", "application/octet-stream"),
            },
        )

        assert resp.status_code == 422

    @pytest.mark.asyncio
    async def test_upload_empty_file(self, client: AsyncClient) -> None:
        """上传空文件应返回 422 错误。"""
        resp = await client.post(
            "/api/v1/upload/document",
            files={
                "file": ("empty.txt", b"", "text/plain"),
            },
        )

        assert resp.status_code == 422

    @pytest.mark.asyncio
    async def test_upload_creates_project(
        self,
        client: AsyncClient,
        db_session: AsyncSession,
    ) -> None:
        """上传应创建新 Project。"""
        pptx_bytes = _make_pptx_bytes()

        resp = await client.post(
            "/api/v1/upload/document",
            files={
                "file": ("math.pptx", pptx_bytes, PPTX_MIME),
            },
        )

        data = resp.json()
        project_id = data["data"]["project_id"]

        # 验证项目已创建
        project_resp = await client.get(f"/api/v1/projects/{project_id}")
        assert project_resp.status_code == 200
        project_data = project_resp.json()
        assert project_data["title"] == "math"

    @pytest.mark.asyncio
    async def test_upload_creates_task(self, client: AsyncClient) -> None:
        """上传应创建新 Task。"""
        pptx_bytes = _make_pptx_bytes()

        resp = await client.post(
            "/api/v1/upload/document",
            files={
                "file": ("lesson.pptx", pptx_bytes, PPTX_MIME),
            },
        )

        data = resp.json()
        task_id = data["data"]["task_id"]

        # 验证任务已创建
        task_resp = await client.get(f"/api/v1/tasks/{task_id}")
        assert task_resp.status_code == 200

    @pytest.mark.asyncio
    async def test_upload_returns_file_info(self, client: AsyncClient) -> None:
        """上传应返回文件信息。"""
        txt_bytes = _make_txt_bytes()

        resp = await client.post(
            "/api/v1/upload/document",
            files={
                "file": ("notes.txt", txt_bytes, "text/plain"),
            },
        )

        data = resp.json()["data"]
        assert data["filename"] == "notes.txt"
        assert data["file_size"] == len(txt_bytes)
        assert data["file_type"] == ".txt"
        assert data["file_path"]  # 非空路径


# ── 脚本 API 测试 ────────────────────────────────────────


class TestScriptsAPI:
    """脚本 (IR) 管理 API 测试。"""

    @pytest.mark.asyncio
    async def test_get_script_not_found(self, client: AsyncClient) -> None:
        """获取不存在的脚本应返回 404。"""
        # 先创建一个项目
        create_resp = await client.post(
            "/api/v1/projects/",
            json={"title": "无脚本项目"},
        )
        project_id = create_resp.json()["id"]

        resp = await client.get(f"/api/v1/scripts/projects/{project_id}/script")
        assert resp.status_code == 404

    @pytest.mark.asyncio
    async def test_generate_script_without_ir(self, client: AsyncClient) -> None:
        """在无 IR 的项目上触发编排应返回 404。"""
        create_resp = await client.post(
            "/api/v1/projects/",
            json={"title": "无IR项目"},
        )
        project_id = create_resp.json()["id"]

        resp = await client.post(
            f"/api/v1/scripts/projects/{project_id}/script/generate"
        )
        assert resp.status_code == 404

    @pytest.mark.asyncio
    async def test_approve_script(self, client: AsyncClient) -> None:
        """审核通过脚本（当前为 stub）。"""
        create_resp = await client.post(
            "/api/v1/projects/",
            json={"title": "待审核"},
        )
        project_id = create_resp.json()["id"]

        resp = await client.post(
            f"/api/v1/scripts/projects/{project_id}/script/approve"
        )
        assert resp.status_code == 200
        assert "审核通过" in resp.json()["message"]
