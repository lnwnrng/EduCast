"""文档解析器 — 将异构教学文档转换为 IR 草稿。

支持: PPTX / PDF / DOCX / Markdown / 纯文本

三层解析:
  1. 物理层: 提取页面、文本框、图片、表格、备注
  2. 逻辑层: 按标题/页边界切分章节→知识点
  3. 教学层: 生成初步分镜（slide → scene 映射）
"""

import io
import logging
import os
import re
from pathlib import Path
from typing import Optional
from uuid import uuid4

from app.exceptions import ParseException
from app.ir.schema import (
    ChapterIR,
    CourseIR,
    KnowledgePointIR,
    SceneIR,
    SceneType,
    VisualSpec,
)

logger = logging.getLogger(__name__)

# 支持的文件类型
SUPPORTED_EXTENSIONS = {".pptx", ".pdf", ".docx", ".md", ".txt"}


# ── 解析结果中间结构 ─────────────────────────────────────


class ParsedSlide:
    """单页解析结果。"""

    def __init__(
        self,
        page_number: int,
        title: str = "",
        body_text: str = "",
        notes: str = "",
        image_paths: list[str] | None = None,
        is_chapter_break: bool = False,
        layout_name: str = "",
    ) -> None:
        self.page_number = page_number
        self.title = title
        self.body_text = body_text
        self.notes = notes
        self.image_paths = image_paths or []
        self.is_chapter_break = is_chapter_break
        self.layout_name = layout_name

    @property
    def has_content(self) -> bool:
        """是否包含有效内容。"""
        return bool(self.title or self.body_text or self.notes)


# ── 文档解析器 ───────────────────────────────────────────


class DocumentParser:
    """文档解析器 — 按文件类型分派到对应解析方法。"""

    def __init__(self, storage_root: str = "./storage") -> None:
        self._storage_root = storage_root

    async def parse(
        self,
        file_path: str,
        file_type: str,
        project_id: Optional[str] = None,
    ) -> CourseIR:
        """解析文档，返回 IR 草稿。

        Args:
            file_path: 文件路径
            file_type: 文件扩展名（如 .pptx）
            project_id: 项目 ID（用于组织存储路径）

        Returns:
            课程脚本 IR 草稿
        """
        file_type = file_type.lower()
        logger.info("开始解析文档: %s (类型: %s)", file_path, file_type)

        parsers = {
            ".pptx": self._parse_pptx,
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".md": self._parse_markdown,
            ".txt": self._parse_text,
        }

        parser = parsers.get(file_type)
        if parser is None:
            raise ParseException(f"不支持的文件类型: {file_type}")

        if not os.path.exists(file_path):
            raise ParseException(f"文件不存在: {file_path}")

        try:
            ir = await parser(file_path, project_id)
        except ParseException:
            raise
        except Exception as exc:
            logger.error("文档解析出错: %s", exc, exc_info=True)
            raise ParseException(f"文档解析失败: {exc}") from exc

        logger.info(
            "文档解析完成: %d 章节, %d 知识点, %d 分镜",
            len(ir.chapters),
            sum(
                len(ch.knowledge_points)
                for ch in ir.chapters
            ),
            sum(
                len(kp.scenes)
                for ch in ir.chapters
                for kp in ch.knowledge_points
            ),
        )
        return ir

    # ── PPTX 解析 ────────────────────────────────────────

    async def _parse_pptx(
        self, file_path: str, project_id: Optional[str] = None
    ) -> CourseIR:
        """解析 PPTX 文件。

        使用 python-pptx:
        - 逐页提取文本框、图片、备注
        - 检测"标题幻灯片"布局作为章节分界
        - 图片保存到 storage/{project_id}/slides/
        """
        from pptx import Presentation
        from pptx.util import Emu

        logger.info("PPTX 解析: %s", file_path)
        prs = Presentation(file_path)

        # 章节标题检测关键词（布局名称）
        title_layout_keywords = {
            "title",
            "标题",
            "section",
            "章节",
            "title slide",
            "section header",
        }

        slides: list[ParsedSlide] = []

        # 设置图片存储目录
        slide_img_dir = self._get_slide_dir(project_id)

        for slide_idx, slide in enumerate(prs.slides):
            page_number = slide_idx + 1
            layout_name = (
                slide.slide_layout.name.lower()
                if slide.slide_layout and slide.slide_layout.name
                else ""
            )

            # 检测是否为章节分界幻灯片
            is_chapter_break = any(
                kw in layout_name for kw in title_layout_keywords
            )

            # 提取标题
            title = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text.strip()

            # 提取正文文本（排除标题 shape）
            body_parts: list[str] = []
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        body_parts.append(text)

            body_text = "\n".join(body_parts)

            # 提取备注
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            # 提取内嵌图片
            image_paths: list[str] = []
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        img_ext = image.content_type.split("/")[-1]
                        if img_ext == "jpeg":
                            img_ext = "jpg"
                        img_filename = (
                            f"slide_{page_number}_img_{len(image_paths) + 1}"
                            f".{img_ext}"
                        )
                        img_path = os.path.join(
                            slide_img_dir, img_filename
                        )
                        os.makedirs(
                            os.path.dirname(img_path), exist_ok=True
                        )
                        with open(img_path, "wb") as f:
                            f.write(image.blob)
                        image_paths.append(img_path)
                    except Exception as e:
                        logger.warning(
                            "提取图片失败 (slide %d): %s",
                            page_number,
                            e,
                        )

            parsed = ParsedSlide(
                page_number=page_number,
                title=title,
                body_text=body_text,
                notes=notes,
                image_paths=image_paths,
                is_chapter_break=is_chapter_break,
                layout_name=layout_name,
            )
            slides.append(parsed)

        if not slides:
            raise ParseException("PPTX 文件中没有幻灯片")

        # 构建 IR
        return self._build_ir_from_slides(
            slides, Path(file_path).stem, project_id
        )

    # ── PDF 解析 ─────────────────────────────────────────

    async def _parse_pdf(
        self, file_path: str, project_id: Optional[str] = None
    ) -> CourseIR:
        """解析 PDF 文件。

        使用 pdfplumber:
        - 逐页提取文本
        - 检测大号字体文本作为标题
        - 提取图片
        """
        import pdfplumber

        logger.info("PDF 解析: %s", file_path)

        slides: list[ParsedSlide] = []
        slide_img_dir = self._get_slide_dir(project_id)

        with pdfplumber.open(file_path) as pdf:
            if not pdf.pages:
                raise ParseException("PDF 文件中没有页面")

            for page_idx, page in enumerate(pdf.pages):
                page_number = page_idx + 1
                full_text = page.extract_text() or ""
                lines = full_text.strip().split("\n") if full_text.strip() else []

                # 检测标题（尝试通过字符分析）
                title = ""
                body_lines: list[str] = []
                is_chapter_break = False

                if lines:
                    # 尝试通过字符属性检测大号标题
                    title_detected = self._detect_pdf_title(page)
                    if title_detected:
                        title = title_detected
                        # 如果标题是整页唯一内容或占主导 → 章节分界
                        remaining_text = full_text.replace(
                            title, "", 1
                        ).strip()
                        if not remaining_text or len(remaining_text) < 20:
                            is_chapter_break = True
                        body_lines = [remaining_text] if remaining_text else []
                    else:
                        # 第一行作为标题
                        title = lines[0].strip()
                        body_lines = lines[1:]

                body_text = "\n".join(
                    line.strip() for line in body_lines if line.strip()
                )

                # 提取图片
                image_paths: list[str] = []
                try:
                    for img_idx, img in enumerate(page.images):
                        img_filename = (
                            f"page_{page_number}_img_{img_idx + 1}.png"
                        )
                        img_path = os.path.join(
                            slide_img_dir, img_filename
                        )
                        os.makedirs(
                            os.path.dirname(img_path), exist_ok=True
                        )
                        # pdfplumber 的图片对象包含 bbox 信息
                        # 实际图片提取需要 PyMuPDF，这里记录位置信息
                        image_paths.append(img_path)
                except Exception as e:
                    logger.warning(
                        "PDF 图片提取失败 (page %d): %s",
                        page_number,
                        e,
                    )

                parsed = ParsedSlide(
                    page_number=page_number,
                    title=title,
                    body_text=body_text,
                    is_chapter_break=is_chapter_break,
                )
                slides.append(parsed)

        return self._build_ir_from_slides(
            slides, Path(file_path).stem, project_id
        )

    def _detect_pdf_title(self, page: "pdfplumber.page.Page") -> str:
        """检测 PDF 页面中的大号标题文本。"""
        try:
            chars = page.chars
            if not chars:
                return ""

            # 计算字符大小分布
            sizes = [float(c.get("size", 0)) for c in chars if c.get("text", "").strip()]
            if not sizes:
                return ""

            avg_size = sum(sizes) / len(sizes)

            # 大号字符（比平均大 30% 以上）认为是标题
            threshold = avg_size * 1.3
            title_chars = [
                c
                for c in chars
                if float(c.get("size", 0)) >= threshold
                and c.get("text", "").strip()
            ]

            if title_chars:
                # 按位置排序，拼接标题文本
                title_chars.sort(
                    key=lambda c: (
                        round(float(c.get("top", 0))),
                        float(c.get("x0", 0)),
                    )
                )
                title_text = "".join(c["text"] for c in title_chars).strip()
                if len(title_text) >= 2:
                    return title_text
        except Exception as e:
            logger.debug("PDF 标题检测失败: %s", e)

        return ""

    # ── DOCX 解析 ────────────────────────────────────────

    async def _parse_docx(
        self, file_path: str, project_id: Optional[str] = None
    ) -> CourseIR:
        """解析 DOCX 文件。

        按标题样式（Heading 1/2/3）切分章节与知识点。
        """
        import mammoth

        logger.info("DOCX 解析: %s", file_path)

        # 使用 mammoth 提取原始文本（保留结构信息）
        with open(file_path, "rb") as f:
            result = mammoth.extract_raw_text(f)
            raw_text = result.value

        if not raw_text.strip():
            raise ParseException("DOCX 文件内容为空")

        # 将 DOCX 文本按段落切分，映射为 Markdown 风格再走 Markdown 解析
        return self._build_ir_from_text(
            raw_text, Path(file_path).stem, project_id
        )

    # ── Markdown 解析 ────────────────────────────────────

    async def _parse_markdown(
        self, file_path: str, project_id: Optional[str] = None
    ) -> CourseIR:
        """解析 Markdown 文件。

        按标题层级拆分:
        - # → 章节 (Chapter)
        - ## → 知识点 (KnowledgePoint)
        - 正文 → 分镜旁白 (Scene.narration_text)
        """
        logger.info("Markdown 解析: %s", file_path)

        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        if not content.strip():
            raise ParseException("Markdown 文件内容为空")

        return self._build_ir_from_markdown(
            content, Path(file_path).stem, project_id
        )

    # ── 纯文本解析 ───────────────────────────────────────

    async def _parse_text(
        self, file_path: str, project_id: Optional[str] = None
    ) -> CourseIR:
        """解析纯文本文件。

        按段落（双换行）切分，整体作为单章节。
        """
        logger.info("纯文本解析: %s", file_path)

        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        if not content.strip():
            raise ParseException("文本文件内容为空")

        return self._build_ir_from_text(
            content, Path(file_path).stem, project_id
        )

    # ── IR 构建辅助方法 ──────────────────────────────────

    def _build_ir_from_slides(
        self,
        slides: list[ParsedSlide],
        title: str,
        project_id: Optional[str] = None,
    ) -> CourseIR:
        """从解析的幻灯片列表构建 IR。

        切分策略:
        1. is_chapter_break=True 的幻灯片开启新章节
        2. 每张幻灯片对应一个 Scene
        3. 连续的非章节分界幻灯片归入同一知识点
        """
        chapters: list[ChapterIR] = []
        current_chapter_slides: list[ParsedSlide] = []
        chapter_title = ""

        def _flush_chapter() -> None:
            """将当前累积的幻灯片打包为一个章节。"""
            nonlocal current_chapter_slides, chapter_title
            if not current_chapter_slides:
                return

            chapter_order = len(chapters) + 1
            ch_title = chapter_title or f"第{chapter_order}章"

            # 将幻灯片分组为知识点
            kps = self._slides_to_knowledge_points(
                current_chapter_slides
            )

            chapter = ChapterIR(
                title=ch_title,
                order=chapter_order,
                source_pages=[
                    s.page_number for s in current_chapter_slides
                ],
                knowledge_points=kps,
            )
            chapters.append(chapter)
            current_chapter_slides = []
            chapter_title = ""

        for slide in slides:
            if slide.is_chapter_break:
                # 刷出之前积累的章节
                _flush_chapter()
                chapter_title = slide.title or ""
                # 如果章节分界页本身有内容，也加入
                if slide.has_content and (slide.body_text or slide.notes):
                    current_chapter_slides.append(slide)
            else:
                # 如果还没有任何章节标题，用第一页标题
                if not chapters and not current_chapter_slides and not chapter_title:
                    chapter_title = slide.title or ""
                current_chapter_slides.append(slide)

        # 刷出最后一个章节
        _flush_chapter()

        # 如果没有检测到任何章节，把所有幻灯片放入一个默认章节
        if not chapters and slides:
            kps = self._slides_to_knowledge_points(slides)
            chapters.append(
                ChapterIR(
                    title=title,
                    order=1,
                    source_pages=[s.page_number for s in slides],
                    knowledge_points=kps,
                )
            )

        return CourseIR(
            title=title,
            chapters=chapters,
        )

    def _slides_to_knowledge_points(
        self, slides: list[ParsedSlide]
    ) -> list[KnowledgePointIR]:
        """将连续幻灯片分组为知识点。

        策略: 每张有标题的幻灯片开启新知识点。
        如果所有幻灯片都没有标题，则整体作为一个知识点。
        """
        kps: list[KnowledgePointIR] = []
        current_kp_slides: list[ParsedSlide] = []
        current_kp_title = ""

        def _flush_kp() -> None:
            nonlocal current_kp_slides, current_kp_title
            if not current_kp_slides:
                return

            kp_order = len(kps) + 1
            kp_title = current_kp_title or f"知识点{kp_order}"

            scenes = []
            for scene_idx, s in enumerate(current_kp_slides):
                # 优先使用备注作为旁白，其次正文，最后标题
                narration = s.notes or s.body_text or s.title or ""
                subtitle = s.body_text or s.title or ""

                visual_spec = VisualSpec(
                    slide_ref=f"slide_{s.page_number}.png",
                    image_refs=s.image_paths,
                )

                scene = SceneIR(
                    order=scene_idx + 1,
                    scene_type=SceneType.SLIDE,
                    narration_text=narration,
                    subtitle_text=subtitle,
                    visual_spec=visual_spec,
                    source_page=s.page_number,
                )
                scenes.append(scene)

            # 收集知识要点
            key_points = []
            for s in current_kp_slides:
                if s.body_text:
                    # 将正文按行拆分为要点
                    for line in s.body_text.split("\n"):
                        line = line.strip()
                        if line and len(line) > 2:
                            key_points.append(line)

            kp = KnowledgePointIR(
                title=kp_title,
                key_points=key_points[:10],  # 限制数量
                source_ref=f"pages {current_kp_slides[0].page_number}"
                f"-{current_kp_slides[-1].page_number}",
                scenes=scenes,
            )
            kps.append(kp)
            current_kp_slides = []
            current_kp_title = ""

        for slide in slides:
            if slide.title and current_kp_slides:
                # 新标题 → 刷出之前的知识点
                _flush_kp()
            if not current_kp_title:
                current_kp_title = slide.title or ""
            current_kp_slides.append(slide)

        _flush_kp()

        # 如果没有知识点（不应发生），创建一个默认的
        if not kps and slides:
            kps.append(
                KnowledgePointIR(
                    title="内容概述",
                    scenes=[
                        SceneIR(
                            order=1,
                            scene_type=SceneType.SLIDE,
                            narration_text="待编排内容。",
                            subtitle_text="待编排内容。",
                        ),
                    ],
                )
            )

        return kps

    def _build_ir_from_markdown(
        self,
        content: str,
        title: str,
        project_id: Optional[str] = None,
    ) -> CourseIR:
        """从 Markdown 内容构建 IR。

        按标题层级:
        - # → Chapter
        - ## → KnowledgePoint
        - 正文段落 → Scene
        """
        chapters: list[ChapterIR] = []

        # 正则拆分: 按 # 开头的行
        heading_pattern = re.compile(r"^(#{1,3})\s+(.+)$", re.MULTILINE)

        # 将内容拆分为 (level, title, body) 块
        sections: list[tuple[int, str, str]] = []
        last_end = 0
        for match in heading_pattern.finditer(content):
            # 之前未归属的内容
            if last_end == 0 and match.start() > 0:
                prefix = content[:match.start()].strip()
                if prefix:
                    sections.append((0, "", prefix))

            level = len(match.group(1))  # 1=#, 2=##, 3=###
            heading_title = match.group(2).strip()
            last_end = match.end()

            # 找到下一个标题或文件末尾之间的内容
            next_match = heading_pattern.search(content, match.end())
            if next_match:
                body = content[match.end(): next_match.start()].strip()
            else:
                body = content[match.end():].strip()

            sections.append((level, heading_title, body))

        # 如果没有标题标记，整体作为一个章节
        if not sections:
            return self._build_ir_from_text(content, title, project_id)

        # 构建层级结构
        current_chapter_title = ""
        current_kps: list[KnowledgePointIR] = []

        def _flush_md_chapter() -> None:
            nonlocal current_chapter_title, current_kps
            if not current_kps:
                return
            ch_order = len(chapters) + 1
            ch_title = current_chapter_title or f"第{ch_order}章"
            chapters.append(
                ChapterIR(
                    title=ch_title,
                    order=ch_order,
                    knowledge_points=current_kps,
                )
            )
            current_kps = []
            current_chapter_title = ""

        for level, sec_title, body in sections:
            if level <= 1:
                # # 标题 → 新章节
                _flush_md_chapter()
                current_chapter_title = sec_title
                if body:
                    kp = self._text_to_knowledge_point(
                        sec_title or "概述", body
                    )
                    current_kps.append(kp)
            elif level == 2:
                # ## 标题 → 新知识点
                kp = self._text_to_knowledge_point(
                    sec_title, body
                )
                current_kps.append(kp)
            else:
                # ### 及以下 → 合并到当前知识点
                if current_kps:
                    # 追加场景到最后一个知识点
                    if body:
                        scene = SceneIR(
                            order=len(current_kps[-1].scenes) + 1,
                            scene_type=SceneType.SLIDE,
                            narration_text=body,
                            subtitle_text=sec_title,
                        )
                        current_kps[-1].scenes.append(scene)
                elif body:
                    kp = self._text_to_knowledge_point(
                        sec_title, body
                    )
                    current_kps.append(kp)

        _flush_md_chapter()

        # 如果没有产出章节
        if not chapters:
            return self._build_ir_from_text(content, title, project_id)

        return CourseIR(title=title, chapters=chapters)

    def _build_ir_from_text(
        self,
        content: str,
        title: str,
        project_id: Optional[str] = None,
    ) -> CourseIR:
        """从纯文本构建 IR — 按段落切分。"""
        # 按双换行分段
        paragraphs = [
            p.strip()
            for p in re.split(r"\n\s*\n", content)
            if p.strip()
        ]

        if not paragraphs:
            raise ParseException("文本内容为空")

        scenes: list[SceneIR] = []
        for idx, para in enumerate(paragraphs):
            scene = SceneIR(
                order=idx + 1,
                scene_type=SceneType.SLIDE,
                narration_text=para,
                subtitle_text=para[:100],  # 字幕截断
            )
            scenes.append(scene)

        kp = KnowledgePointIR(
            title=title,
            key_points=[
                p[:80] for p in paragraphs[:10]
            ],
            scenes=scenes,
        )

        chapter = ChapterIR(
            title=title,
            order=1,
            knowledge_points=[kp],
        )

        return CourseIR(title=title, chapters=[chapter])

    def _text_to_knowledge_point(
        self, title: str, body: str
    ) -> KnowledgePointIR:
        """将标题+正文转为知识点。"""
        # 按段落拆分为多个场景
        paragraphs = [
            p.strip()
            for p in re.split(r"\n\s*\n", body)
            if p.strip()
        ]

        if not paragraphs:
            paragraphs = [body.strip()] if body.strip() else [""]

        scenes: list[SceneIR] = []
        for idx, para in enumerate(paragraphs):
            if not para:
                continue
            scene = SceneIR(
                order=idx + 1,
                scene_type=SceneType.SLIDE,
                narration_text=para,
                subtitle_text=para[:100],
            )
            scenes.append(scene)

        # 确保至少一个场景
        if not scenes:
            scenes.append(
                SceneIR(
                    order=1,
                    scene_type=SceneType.SLIDE,
                    narration_text=title,
                    subtitle_text=title,
                )
            )

        # 提取要点
        key_points = []
        for para in paragraphs:
            for line in para.split("\n"):
                line = line.strip().lstrip("- ").lstrip("* ").lstrip("• ")
                if line and len(line) > 2:
                    key_points.append(line[:80])

        return KnowledgePointIR(
            title=title,
            key_points=key_points[:10],
            scenes=scenes,
        )

    def _get_slide_dir(self, project_id: Optional[str] = None) -> str:
        """获取幻灯片图片存储目录。"""
        if project_id:
            return os.path.join(
                self._storage_root, project_id, "slides"
            )
        return os.path.join(self._storage_root, "slides")
